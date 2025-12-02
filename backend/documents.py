"""
Document handling module for LLM Council.

Handles document upload, storage, text extraction, and management.
Extracted text is stored separately from original files for efficient retrieval.
"""

import os
import json
import uuid
from datetime import datetime
from typing import Dict, List, Optional
from pathlib import Path

# Constants
DOCUMENTS_DIR = "data/documents"
REGISTRY_FILE = "data/document_registry.json"
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
MAX_TEXT_LENGTH = 500 * 1024  # 500KB

SUPPORTED_EXTENSIONS = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
}


def ensure_documents_dir() -> None:
    """Create data/documents directory if it doesn't exist."""
    os.makedirs(DOCUMENTS_DIR, exist_ok=True)


def load_registry() -> Dict[str, Dict]:
    """
    Load document registry from JSON file.

    Returns:
        Dict mapping document IDs to metadata entries.
    """
    if not os.path.exists(REGISTRY_FILE):
        return {}

    try:
        with open(REGISTRY_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    except (json.JSONDecodeError, IOError) as e:
        print(f"Error loading registry: {e}")
        return {}


def save_registry(registry: Dict[str, Dict]) -> None:
    """
    Save document registry to JSON file.

    Args:
        registry: Dict mapping document IDs to metadata entries.
    """
    try:
        # Ensure parent directory exists
        os.makedirs(os.path.dirname(REGISTRY_FILE), exist_ok=True)

        with open(REGISTRY_FILE, 'w', encoding='utf-8') as f:
            json.dump(registry, f, indent=2, ensure_ascii=False)
    except IOError as e:
        print(f"Error saving registry: {e}")
        raise


def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from PDF file using PyPDF2.

    Args:
        file_path: Path to PDF file.

    Returns:
        Extracted text content.
    """
    try:
        import PyPDF2

        text_parts = []
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")
                except Exception as e:
                    text_parts.append(f"--- Page {page_num + 1} ---\n[Error extracting page: {e}]")

        return "\n\n".join(text_parts) if text_parts else "[No text content extracted from PDF]"

    except ImportError:
        return "[Error: PyPDF2 not installed. Install with: pip install PyPDF2]"
    except Exception as e:
        return f"[Error extracting PDF: {str(e)}]"


def extract_text_from_docx(file_path: str) -> str:
    """
    Extract text from DOCX file using python-docx.

    Args:
        file_path: Path to DOCX file.

    Returns:
        Extracted text content.
    """
    try:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

        return "\n\n".join(paragraphs) if paragraphs else "[No text content extracted from DOCX]"

    except ImportError:
        return "[Error: python-docx not installed. Install with: pip install python-docx]"
    except Exception as e:
        return f"[Error extracting DOCX: {str(e)}]"


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from PPTX file using python-pptx.

    Args:
        file_path: Path to PPTX file.

    Returns:
        Extracted text content.
    """
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))

        return "\n\n".join(text_parts) if text_parts else "[No text content extracted from PPTX]"

    except ImportError:
        return "[Error: python-pptx not installed. Install with: pip install python-pptx]"
    except Exception as e:
        return f"[Error extracting PPTX: {str(e)}]"


def extract_text_from_txt(file_path: str) -> str:
    """
    Extract text from plain text or markdown file.

    Args:
        file_path: Path to text file.

    Returns:
        Text content.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encoding
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            return f"[Error reading text file: {str(e)}]"
    except Exception as e:
        return f"[Error reading text file: {str(e)}]"


def extract_text(file_path: str, extension: str) -> str:
    """
    Route text extraction to appropriate handler based on file extension.

    Args:
        file_path: Path to file.
        extension: File extension (e.g., '.pdf').

    Returns:
        Extracted text content or metadata for images.
    """
    extension = extension.lower()

    # Image files - return metadata only
    if extension in ['.png', '.jpg', '.jpeg', '.gif', '.webp']:
        file_size = os.path.getsize(file_path)
        return f"[Image file: {os.path.basename(file_path)} - {file_size} bytes. No OCR performed.]"

    # Text extraction by file type
    if extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif extension == '.docx':
        return extract_text_from_docx(file_path)
    elif extension == '.pptx':
        return extract_text_from_pptx(file_path)
    elif extension in ['.txt', '.md']:
        return extract_text_from_txt(file_path)
    else:
        return f"[Unsupported file type: {extension}]"


async def save_document(file_content: bytes, filename: str) -> Dict:
    """
    Save uploaded document, extract text, and create registry entry.

    Args:
        file_content: Raw file bytes.
        filename: Original filename.

    Returns:
        Document metadata dict.

    Raises:
        ValueError: If file is invalid or too large.
    """
    # Validate file size
    if len(file_content) > MAX_FILE_SIZE:
        raise ValueError(f"File too large. Maximum size is {MAX_FILE_SIZE / (1024*1024):.1f}MB")

    # Validate extension
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {extension}")

    # Ensure directories exist
    ensure_documents_dir()

    # Generate unique ID and paths
    doc_id = str(uuid.uuid4())
    original_file_path = os.path.join(DOCUMENTS_DIR, f"{doc_id}{extension}")
    text_file_path = os.path.join(DOCUMENTS_DIR, f"{doc_id}.txt")

    # Save original file
    try:
        with open(original_file_path, 'wb') as f:
            f.write(file_content)
    except IOError as e:
        raise ValueError(f"Failed to save file: {str(e)}")

    # Extract text
    extracted_text = extract_text(original_file_path, extension)

    # Truncate if too long
    text_truncated = False
    if len(extracted_text) > MAX_TEXT_LENGTH:
        extracted_text = extracted_text[:MAX_TEXT_LENGTH] + "\n\n[... Text truncated due to length ...]"
        text_truncated = True

    # Save extracted text to separate file
    try:
        with open(text_file_path, 'w', encoding='utf-8') as f:
            f.write(extracted_text)
    except IOError as e:
        # Clean up original file if text save fails
        os.remove(original_file_path)
        raise ValueError(f"Failed to save extracted text: {str(e)}")

    # Create registry entry (metadata only)
    metadata = {
        "id": doc_id,
        "filename": filename,
        "extension": extension,
        "size": len(file_content),
        "uploaded_at": datetime.utcnow().isoformat(),
        "text_length": len(extracted_text),
        "text_truncated": text_truncated,
        "is_active": True
    }

    # Update registry
    registry = load_registry()
    registry[doc_id] = metadata
    save_registry(registry)

    return metadata


def get_document(doc_id: str) -> Optional[Dict]:
    """
    Get document metadata from registry.

    Args:
        doc_id: Document ID.

    Returns:
        Document metadata dict or None if not found.
    """
    registry = load_registry()
    return registry.get(doc_id)


def get_document_text(doc_id: str) -> Optional[str]:
    """
    Read extracted text from document's .txt file.

    Args:
        doc_id: Document ID.

    Returns:
        Extracted text content or None if not found.
    """
    text_file_path = os.path.join(DOCUMENTS_DIR, f"{doc_id}.txt")

    if not os.path.exists(text_file_path):
        return None

    try:
        with open(text_file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"Error reading document text: {e}")
        return None


def list_documents() -> List[Dict]:
    """
    List all documents with metadata and text preview.

    Returns:
        List of document metadata dicts with preview field added.
    """
    registry = load_registry()
    documents = []

    for doc_id, metadata in registry.items():
        doc_with_preview = metadata.copy()

        # Add text preview (first 200 chars)
        text = get_document_text(doc_id)
        if text:
            preview = text[:200].strip()
            if len(text) > 200:
                preview += "..."
            doc_with_preview["preview"] = preview
        else:
            doc_with_preview["preview"] = "[No preview available]"

        documents.append(doc_with_preview)

    # Sort by upload date (newest first)
    documents.sort(key=lambda x: x["uploaded_at"], reverse=True)

    return documents


def delete_document(doc_id: str) -> bool:
    """
    Delete document files and remove from registry.

    Args:
        doc_id: Document ID.

    Returns:
        True if successful, False if document not found.
    """
    registry = load_registry()

    if doc_id not in registry:
        return False

    metadata = registry[doc_id]
    extension = metadata["extension"]

    # Delete original file
    original_file_path = os.path.join(DOCUMENTS_DIR, f"{doc_id}{extension}")
    if os.path.exists(original_file_path):
        try:
            os.remove(original_file_path)
        except Exception as e:
            print(f"Error deleting original file: {e}")

    # Delete text file
    text_file_path = os.path.join(DOCUMENTS_DIR, f"{doc_id}.txt")
    if os.path.exists(text_file_path):
        try:
            os.remove(text_file_path)
        except Exception as e:
            print(f"Error deleting text file: {e}")

    # Remove from registry
    del registry[doc_id]
    save_registry(registry)

    return True


def toggle_document_active(doc_id: str, is_active: bool) -> bool:
    """
    Toggle document active status for context inclusion.

    Args:
        doc_id: Document ID.
        is_active: New active status.

    Returns:
        True if successful, False if document not found.
    """
    registry = load_registry()

    if doc_id not in registry:
        return False

    registry[doc_id]["is_active"] = is_active
    save_registry(registry)

    return True


def get_active_documents_context() -> str:
    """
    Build context string from all active documents.

    Returns:
        Formatted string with all active document contents.
    """
    registry = load_registry()
    active_docs = [
        (doc_id, metadata)
        for doc_id, metadata in registry.items()
        if metadata.get("is_active", True)
    ]

    if not active_docs:
        return ""

    context_parts = ["=== UPLOADED DOCUMENTS ===\n"]

    for doc_id, metadata in active_docs:
        filename = metadata["filename"]
        text = get_document_text(doc_id)

        if text:
            context_parts.append(f"--- Document: {filename} ---")
            context_parts.append(text)
            context_parts.append("")  # Empty line separator

    context_parts.append("=== END DOCUMENTS ===")

    return "\n".join(context_parts)
